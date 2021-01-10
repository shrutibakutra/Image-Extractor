from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageOps
import sys
import datetime
import json
import os

filename= sys.argv[1]

n=0
files = []
comments = []

 # create directory to store extracted images
path = './Images_'+filename
os.makedirs(path)

def write_image(shape):
    global n
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    now = datetime.datetime.now()
    timestamp = str(now.strftime("%Y%m%d_%H-%M-%S"))

    imagepost = str(n) + str(timestamp)   
   

    image_filename = path+'/image{}_{}'.format(imagepost, image.ext)
    image_filename_round = './Images/round_image_' + imagepost +'.png' 
    image_filename_square = './Images/square_image_{}_{}'.format(imagepost, image.ext)    
    image_filename = image_filename + '.'+image.ext
    #  image_filename = 'public/uploads/pages/image{}_{}'.format(n,timestamp+".jpg")    
    files.append({"image": image_filename, "ext":image.ext,"s_image":image_filename_square, "r_image":image_filename_round, "note": "", "isPortrait": 0})
    # store images in perticular location
    with open(image_filename, 'wb') as f:
        f.write(image_bytes)                
    n += 1


def visitor(shape):
    ''' Check a type of an image'''

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape)


def iter_picture_shapes(prs):
    n2 = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            visitor(shape)
        if slide.has_notes_slide:    
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame            
            files[n2]["note"] = text_frame.text.encode('utf-8')
        n2 += 1
    # Check if image is portrait   
    n3 = 0
    for efile in files:
        im = Image.open(efile['image'])
        w = im.size[0]
        h = im.size[1]   

        files[n3]["isPortrait"] = 1 if w < h else 0        
        im1 = im.resize((320, 320))
       
'''API chaining'''
iter_picture_shapes(Presentation(filename))
